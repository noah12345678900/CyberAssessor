"""Generate the executive overview deck for Cybersecurity Assessor.

One-off content builder (python-pptx). Run with the backend venv:
    backend/.venv/Scripts/python.exe build_overview_deck.py
Writes Cybersecurity_Assessor_Overview.pptx next to this script.

Palette is Nuon brand (navy + blue), pulled from the marketing site CSS.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---- Nuon palette (navy / blue accent) ------------------------------------
NAVY = RGBColor(0x0A, 0x25, 0x40)        # --navy
NAVY_DEEP = RGBColor(0x06, 0x1A, 0x30)   # --navy-deep
NAVY2 = RGBColor(0x10, 0x2E, 0x52)       # card-on-dark
NAVY3 = RGBColor(0x0F, 0x33, 0x66)       # gradient mid
BLUE = RGBColor(0x25, 0x63, 0xEB)        # --blue
BLUE_BRIGHT = RGBColor(0x3B, 0x82, 0xF6) # --blue-bright
SKY = RGBColor(0x60, 0xA5, 0xFA)         # light blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE7, 0xEF, 0xFB)       # pale blue text on dark
SOFT = RGBColor(0xF6, 0xF9, 0xFC)        # --bg-soft (card on light)
LINE = RGBColor(0xE5, 0xEB, 0xF3)        # --line
SLATE = RGBColor(0x5B, 0x6B, 0x85)       # --muted
INK = RGBColor(0x1A, 0x25, 0x40)         # --ink
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xE0, 0x8A, 0x1E)
MIST = RGBColor(0xC8, 0xD4, 0xE8)        # footnote on dark

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# primitives
# ---------------------------------------------------------------------------
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _grad_bg(slide, c1, c2, angle=90):
    """Full-bleed gradient rectangle (added first => sits behind content)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    sp.fill.gradient()
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _rect(slide, l, t, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def _grad_rect(slide, l, t, w, h, c1, c2, angle=0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.gradient()
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _round(slide, l, t, w, h, fill, line=None, line_w=None, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1.25)
    sp.shadow.inherit = False
    return sp


def _grad_round(slide, l, t, w, h, c1, c2, angle=45, line=None, line_w=None, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.gradient()
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1.25)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is a list of run tuples.

    A run tuple is (text, size, color, bold, italic) or, to make the run a
    hyperlink, (text, size, color, bold, italic, url).
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for run in para:
            txt, size, color, bold, italic = run[:5]
            url = run[5] if len(run) > 5 else None
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Segoe UI"
            if url:
                r.hyperlink.address = url
    return tb


def _chip(slide, l, t, w, h, label, fill, txtcolor=WHITE, size=11, line=None):
    c = _round(slide, l, t, w, h, fill, line=line, line_w=Pt(1), radius=0.5)
    tf = c.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = txtcolor
    r.font.name = "Segoe UI"
    return c


def _accentbar(slide):
    _grad_rect(slide, 0, 0, W, Inches(0.14), BLUE, BLUE_BRIGHT, angle=0)


def _kicker(slide, l, t, label, color=BLUE):
    _text(slide, l, t, Inches(6), Inches(0.35),
          [[(label.upper(), 12, color, True, False)]])


# ===========================================================================
# Slide 1 — Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_grad_bg(s, NAVY_DEEP, NAVY3, angle=60)
_grad_rect(s, 0, 0, W, Inches(0.16), BLUE_BRIGHT, BLUE, angle=0)
_rect(s, 0, H - Inches(0.1), W, Inches(0.1), BLUE)

# soft glow accent circle (fully on-slide, low-contrast)
_grad_round(s, Inches(8.75), Inches(1.55), Inches(4.4), Inches(4.4),
            NAVY2, NAVY_DEEP, angle=45, radius=0.5)

# accent rail
_grad_round(s, Inches(0.9), Inches(1.95), Inches(0.18), Inches(2.35),
            BLUE_BRIGHT, BLUE, angle=90, radius=0.5)

_kicker(s, Inches(1.32), Inches(1.55), "Nuon  ·  AI-native control assessment")
_text(s, Inches(1.3), Inches(1.95), Inches(11), Inches(1.4),
      [[("Cybersecurity Assessor", 48, WHITE, True, False)]])
_text(s, Inches(1.32), Inches(3.0), Inches(10.6), Inches(1.0),
      [[("Reasons like a senior assessor, verifies like an auditor — backing every verdict with a cited artifact.",
         22, SKY, False, False)]])
_text(s, Inches(1.32), Inches(3.95), Inches(10.4), Inches(1.4),
      [[("Turn months of manual control assessment into hours. Defensible verdicts, cited",
         16, LIGHT, False, False)],
       [("evidence, and an auditable trail — running on hardware you control, with the AI endpoint you choose.",
         16, LIGHT, False, False)]])

# classification spectrum chips
_chip(s, Inches(1.32), Inches(5.25), Inches(1.9), Inches(0.46), "UNCLASSIFIED", NAVY2, SKY, 10.5, line=BLUE)
_chip(s, Inches(3.34), Inches(5.25), Inches(1.1), Inches(0.46), "CUI", NAVY2, SKY, 10.5, line=BLUE)
_chip(s, Inches(4.56), Inches(5.25), Inches(1.4), Inches(0.46), "SECRET", NAVY2, SKY, 10.5, line=BLUE)
_chip(s, Inches(6.08), Inches(5.25), Inches(2.0), Inches(0.46), "TOP SECRET / SCI", BLUE, WHITE, 10.5)

# capability chips
_chip(s, Inches(1.32), Inches(5.95), Inches(2.5), Inches(0.5), "7 frameworks · eMASS-native", NAVY2, LIGHT, 11, line=BLUE)
_chip(s, Inches(3.96), Inches(5.95), Inches(2.5), Inches(0.5), "Runs in-enclave · no SaaS", NAVY2, LIGHT, 11, line=BLUE)

_text(s, Inches(1.3), Inches(6.75), Inches(11), Inches(0.5),
      [[("Self-contained Windows desktop application  ·  v2.0.1", 13, MIST, False, False)]])


# ===========================================================================
# Slide 2 — The problem
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.5), "The problem", AMBER)
_text(s, Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.9),
      [[("Control assessment is the bottleneck to authorization", 32, NAVY, True, False)]])
_text(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.5),
      [[("Every ATO, IATT, and continuous-monitoring cycle hinges on the same slow, manual work.",
         16, SLATE, False, False)]])

pains = [
    ("Hundreds of controls, by hand",
     "Assessors read each control, hunt for the right artifact, and hand-write a justification — for hundreds of CCIs, every cycle."),
    ("Evidence scattered everywhere",
     "Policies, STIG checklists, scans, screenshots, and diagrams live across SharePoint, shares, and inboxes — finding the right one eats hours."),
    ("Findings that don't hold up",
     "Weak citations and inconsistent narratives get bounced by a 3PAO or JAB reviewer, forcing expensive rework late in the process."),
    ("Sensitive data can't leave",
     "CUI and classified program data can't be pasted into a public AI tool — so teams stay manual while the rest of the world automates."),
]
cards_y = Inches(2.35)
cw, gap = Inches(5.75), Inches(0.4)
for i, (h, b) in enumerate(pains):
    col = i % 2
    rowi = i // 2
    x = Inches(0.8) + (cw + gap) * col
    y = cards_y + (Inches(2.0)) * rowi
    _round(s, x, y, cw, Inches(1.78), SOFT, line=LINE, line_w=Pt(1))
    _grad_round(s, x, y, Inches(0.14), Inches(1.78), AMBER, RGBColor(0xF2, 0xB1, 0x55), angle=90, radius=0.5)
    _text(s, x + Inches(0.4), y + Inches(0.22), cw - Inches(0.65), Inches(0.5),
          [[(h, 18, NAVY, True, False)]])
    _text(s, x + Inches(0.4), y + Inches(0.78), cw - Inches(0.65), Inches(0.9),
          [[(b, 13, INK, False, False)]], line_spacing=1.12)


# ===========================================================================
# Slide 3 — Classification & privacy boundary  (NEW)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_grad_bg(s, NAVY_DEEP, NAVY3, angle=60)
_grad_rect(s, 0, 0, W, Inches(0.14), BLUE_BRIGHT, BLUE, angle=0)
_kicker(s, Inches(0.8), Inches(0.45), "Built for every enclave", SKY)
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("From Unclassified to Top Secret — you control your data", 30, WHITE, True, False)]])
_text(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.5),
      [[("The same engine runs from a CUI dev box to a TS/SCI accreditation. Only the AI endpoint changes.",
         15, SKY, False, False)]])

# classification spectrum bar
seg_labels = ["UNCLASSIFIED", "CUI", "SECRET", "TOP SECRET / SCI"]
seg_w = [Inches(3.0), Inches(2.4), Inches(2.9), Inches(3.43)]
sx = Inches(0.8)
sy = Inches(2.25)
sh = Inches(0.62)
seg_fills = [NAVY2, NAVY2, BLUE, BLUE_BRIGHT]
seg_txt = [SKY, SKY, WHITE, WHITE]
for lab, wseg, fill, tc in zip(seg_labels, seg_w, seg_fills, seg_txt):
    seg = _round(s, sx, sy, wseg - Inches(0.1), sh, fill, line=BLUE, line_w=Pt(1), radius=0.5)
    tf = seg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = lab
    r.font.size = Pt(12.5)
    r.font.bold = True
    r.font.color.rgb = tc
    r.font.name = "Segoe UI"
    sx = sx + wseg

# three privacy pillars
priv = [
    ("Lives inside your enclave",
     "Installs and runs entirely on the assessor's workstation — no SaaS, no portal, no account. Your eMASS workbook and every artifact stay on the network they already live on; the engine and your chosen AI endpoint work together from inside your boundary."),
    ("You choose the AI boundary",
     "Point it at the model approved for your environment: a commercial key for unclassified work, or a GovCloud / in-boundary endpoint for higher enclaves. The only outbound call is to the endpoint you designate."),
    ("Nothing phones home",
     "No telemetry, no analytics, no evidence files uploaded. Files stay on the workstation; only the prompt you assess — control text plus the evidence excerpts it cites — goes to the AI endpoint you designate. On-device learning stays local."),
]
py = Inches(3.25)
pw = Inches(3.84)
pgap = Inches(0.1)
for i, (title, body) in enumerate(priv):
    x = Inches(0.8) + (pw + pgap) * i
    _round(s, x, py, pw, Inches(3.2), NAVY2, line=BLUE, line_w=Pt(1), radius=0.06)
    _grad_rect(s, x, py, pw, Inches(0.09), BLUE_BRIGHT, BLUE, angle=0)
    _text(s, x + Inches(0.32), py + Inches(0.28), pw - Inches(0.6), Inches(0.72),
          [[(title, 17, WHITE, True, False)]], line_spacing=1.02)
    _text(s, x + Inches(0.32), py + Inches(1.04), pw - Inches(0.6), Inches(2.0),
          [[(body, 12, LIGHT, False, False)]], line_spacing=1.13)

# ===========================================================================
# Slide 4 — How it works (the pipeline + kernel descriptor)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_grad_bg(s, NAVY_DEEP, NAVY3, angle=60)
_grad_rect(s, 0, 0, W, Inches(0.14), BLUE_BRIGHT, BLUE, angle=0)
_kicker(s, Inches(0.8), Inches(0.4), "How it works", SKY)
_text(s, Inches(0.8), Inches(0.74), Inches(11.7), Inches(0.8),
      [[("AI speed, kernel-grade rigor", 32, WHITE, True, False)]])
_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
      [[("We didn't bolt a chatbot onto a spreadsheet. We built a reasoning engine that mimics how a senior assessor works — then checks the AI's judgment against it.",
         14, SKY, False, False)]])

steps = [
    ("1 · Ingest", "Reads your eMASS CCIS workbook plus every evidence artifact — documents, scans, STIG checklists, network diagrams, and screenshots."),
    ("2 · Correlate", "A proprietary 5-tier evidence engine links each artifact to the exact control it proves — by document number, CCI, control ID, content type, then ML relevance."),
    ("3 · Reason", "Rules modeled on assessor judgment auto-resolve the clear cases — inheritance, scope-exclusion, provider-owned — with no LLM call and no guessing."),
    ("4 · Validate", "The judgment calls go to the LLM, which proposes a cited verdict — then a second deterministic pass validates it before it's ever accepted."),
    ("5 · Write back", "Verdicts, dates, and narratives written straight into the eMASS workbook — formatting, comments, and data validation fully preserved."),
]
bx = Inches(0.8)
bw = Inches(2.3)
bgap = Inches(0.16)
by = Inches(2.25)
bh = Inches(2.45)
for i, (title, body) in enumerate(steps):
    x = bx + (bw + bgap) * i
    _round(s, x, by, bw, bh, NAVY2, line=BLUE, line_w=Pt(1), radius=0.06)
    _grad_rect(s, x, by, bw, Inches(0.5), BLUE, BLUE_BRIGHT, angle=0)
    _text(s, x + Inches(0.14), by + Inches(0.04), bw - Inches(0.28), Inches(0.42),
          [[(title, 15, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x + Inches(0.2), by + Inches(0.62), bw - Inches(0.4), Inches(1.77),
          [[(body, 12, LIGHT, False, False)]], line_spacing=1.12)

# kernel descriptor strip
_round(s, Inches(0.8), Inches(5.0), Inches(11.73), Inches(2.0), NAVY2, line=BLUE, line_w=Pt(1.25), radius=0.06)
_grad_rect(s, Inches(0.8), Inches(5.0), Inches(0.18), Inches(2.0), BLUE_BRIGHT, BLUE, angle=90)
_text(s, Inches(1.18), Inches(5.14), Inches(11.1), Inches(0.5),
      [[("The Assessment Engine — two layers of intelligence, one defensible answer", 18, SKY, True, False)]])
_text(s, Inches(1.18), Inches(5.62), Inches(11.1), Inches(1.32),
      [[("At the core is a proprietary reasoning engine that encodes how an experienced assessor actually works — which evidence matters, when a "
         "control is inherited, when scope excludes it, when the proof simply isn't there. It resolves everything it can prove on its own, then "
         "hands only the true judgment calls to the LLM as a second opinion — and re-validates that opinion before accepting it. ",
         13, LIGHT, False, False)],
       [("The deterministic engine does the heavy lifting for speed and validates the AI's second opinion for trust — and when the evidence is ambiguous, it abstains rather than guess.",
         13, WHITE, True, False)]], line_spacing=1.1)


# ===========================================================================
# Slide 5 — Under the hood: proprietary engine internals
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.45), "Under the hood")
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("Three engines that learn and adapt", 32, NAVY, True, False)]])
_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
      [[("Not off-the-shelf prompting — proprietary intelligence, all running locally.", 16, SLATE, False, False)]])

eng = [
    ("Machine-learning evidence sweep",
     "Boundary-aware SharePoint triage",
     "Point it at a 4,000-file SharePoint site and it surfaces the ~30 files that belong to THIS system — reading only metadata and snippets, never downloading. It scores each candidate and proposes the control it proves before a byte is pulled."),
    ("A model that learns your judgment",
     "Learning",
     "Each time an assessor includes or rejects a swept file, an online ML model shifts its scoring toward that choice. It learns what THIS team treats as relevant and sharpens every assessment — no retraining, nothing leaving the box."),
    ("A lie detector for vendor spreadsheets",
     "Adversarial CRM anomaly detection",
     "When a vendor claims 'we handle this control,' the system can auto-pass it — a real time-saver, IF the sheet is honest. This engine flags 'too good to be true' matrices that contradict your scans, so no control is rubber-stamped."),
]
ey = Inches(2.25)
ew = Inches(3.84)
egap = Inches(0.1)
for i, (title, tag, body) in enumerate(eng):
    x = Inches(0.8) + (ew + egap) * i
    _grad_round(s, x, ey, ew, Inches(4.35), NAVY, NAVY3, angle=60, radius=0.05)
    _grad_rect(s, x, ey, ew, Inches(0.09), BLUE_BRIGHT, BLUE, angle=0)
    _text(s, x + Inches(0.3), ey + Inches(0.28), ew - Inches(0.6), Inches(0.8),
          [[(title, 16, WHITE, True, False)]], line_spacing=1.02)
    _chip(s, x + Inches(0.3), ey + Inches(1.02), ew - Inches(0.6), Inches(0.5), tag, BLUE, WHITE, 13)
    _text(s, x + Inches(0.3), ey + Inches(1.72), ew - Inches(0.6), Inches(2.5),
          [[(body, 13, LIGHT, False, False)]], line_spacing=1.2)
_text(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
      [[("All learning happens locally — the model's training and your team's judgment never leave the workstation.",
         12.5, BLUE, False, True)]])


# ===========================================================================
# Slide 6 — Feature grid
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.45), "What it does")
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("One unified workflow, end to end", 32, NAVY, True, False)]])

feats = [
    ("Multi-framework", "NIST 800-53, 800-171 (CMMC), CSF 2.0, ISO 27001, CIS v8, PCI DSS, and SOC 2 — one engine, seven frameworks."),
    ("Reads every artifact", "PDF, Word, PowerPoint, Excel, STIG .ckl/.cklb/XCCDF, Nessus/ACAS scans, Visio diagrams, and screenshots."),
    ("OCR + Vision", "OCR pulls verbatim text from config screenshots; an AI vision pass also reads diagrams and dashboards OCR can't — so image evidence counts."),
    ("Cited narratives", "Every verdict names the exact document, section, and STIG rule it relies on — no invented citations."),
    ("Multi-boundary aware", "Splits responsibility per cloud tenant (e.g. AWS GovCloud vs. Azure Government) with attribution that can't cross boundaries."),
    ("eMASS round-trip", "Reads the CCIS export and writes results back in place — comments, conditional formatting, and data validation survive."),
    ("Auto POA&M", "Generates Plan of Action & Milestones entries for gaps, grounded in the actual finding and severity-based timelines."),
    ("Full audit trail", "Records the precise evidence snippet behind every decision, so a reviewer can replay exactly what the assessor saw."),
    ("Customer responsibility", "Ingests Customer Responsibility Matrices and auto-resolves inherited / provider / shared controls."),
]
gx, gy = Inches(0.8), Inches(1.7)
gw, gh = Inches(3.84), Inches(1.62)
ggap = Inches(0.1)
for i, (h, b) in enumerate(feats):
    col = i % 3
    rowi = i // 3
    x = gx + (gw + ggap) * col
    y = gy + (gh + ggap) * rowi
    _round(s, x, y, gw, gh, SOFT, line=LINE, line_w=Pt(1), radius=0.06)
    _grad_rect(s, x, y, gw, Inches(0.08), BLUE, BLUE_BRIGHT, angle=0)
    _text(s, x + Inches(0.28), y + Inches(0.2), gw - Inches(0.5), Inches(0.45),
          [[(h, 15.5, NAVY, True, False)]])
    _text(s, x + Inches(0.28), y + Inches(0.64), gw - Inches(0.5), Inches(0.9),
          [[(b, 12, INK, False, False)]], line_spacing=1.1)


# ===========================================================================
# Slide 7 — Why it's different
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.45), "Why it's different")
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("A purpose-built assessment engine", 32, NAVY, True, False)]])
_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
      [[("Not a chatbot bolted onto compliance.", 16, SLATE, False, False)]])

diffs = [
    ("Proprietary reasoning engine",
     "Not prompt-wrapping. A purpose-built engine encodes assessor judgment and does the assessing — calling the LLM only for the hard judgment calls, then checking its answer."),
    ("Learns locally, privately",
     "On-device machine learning adapts to how your team works — and the model's training never leaves the workstation. Files stay local; only the prompt you assess reaches the AI endpoint you choose. OCR + Vision included."),
    ("Traceable by design",
     "Every verdict is rule-validated and traced to specific evidence — the LLM only proposes, the engine decides. Built to survive 3PAO / JAB scrutiny."),
    ("Speaks eMASS natively",
     "Consumes the real CCIS workbook and writes results back in place — no re-keying, no broken formatting, no separate system to maintain."),
]
dy = Inches(2.15)
dw = Inches(5.75)
dgap = Inches(0.4)
for i, (h, b) in enumerate(diffs):
    col = i % 2
    rowi = i // 2
    x = Inches(0.8) + (dw + dgap) * col
    y = dy + Inches(2.15) * rowi
    _grad_round(s, x, y, dw, Inches(1.95), NAVY, NAVY3, angle=60, radius=0.06)
    _grad_round(s, x, y, Inches(0.14), Inches(1.95), BLUE_BRIGHT, BLUE, angle=90, radius=0.5)
    _text(s, x + Inches(0.42), y + Inches(0.25), dw - Inches(0.72), Inches(0.5),
          [[(h, 19, WHITE, True, False)]])
    _text(s, x + Inches(0.42), y + Inches(0.85), dw - Inches(0.72), Inches(1.0),
          [[(b, 13, LIGHT, False, False)]], line_spacing=1.15)


# ===========================================================================
# Slide 8 — ROI / estimated savings per assessment
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.45), "Return on investment")
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("What one assessment is worth", 32, NAVY, True, False)]])

roi = [
    ("$700 / control", 26, "manual A&A benchmark", "≈ $233 per CCI average"),
    ("~$180K", 38, "saved per assessment", "manual A&A labor the tool auto-resolves at that rate"),
    ("~2,080 hrs", 34, "assessor time saved", "≈ 52 assessor work-weeks off the A&A effort"),
    ("~80%", 38, "less hands-on effort", "clear cases auto-resolve; experts review the rest"),
]
rx = Inches(0.8)
rw = Inches(2.85)
rgap = Inches(0.13)
for i, (big, big_sz, mid, small) in enumerate(roi):
    x = rx + (rw + rgap) * i
    _grad_round(s, x, Inches(2.2), rw, Inches(2.45), NAVY, NAVY3, angle=60, radius=0.07)
    _grad_rect(s, x, Inches(2.2), rw, Inches(0.1), BLUE_BRIGHT, BLUE, angle=0)
    _text(s, x, Inches(2.5), rw, Inches(0.85),
          [[(big, big_sz, SKY, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x + Inches(0.2), Inches(3.35), rw - Inches(0.4), Inches(0.5),
          [[(mid, 14, WHITE, True, False)]], align=PP_ALIGN.CENTER)
    _text(s, x + Inches(0.25), Inches(3.85), rw - Inches(0.5), Inches(0.8),
          [[(small, 11.5, LIGHT, False, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)

# methodology / asterisk box
_round(s, Inches(0.8), Inches(4.78), Inches(11.73), Inches(2.55), SOFT, line=LINE, line_w=Pt(1), radius=0.06)
_grad_round(s, Inches(0.8), Inches(4.78), Inches(0.14), Inches(2.55), AMBER, RGBColor(0xF2, 0xB1, 0x55), angle=90, radius=0.5)
_text(s, Inches(1.15), Inches(4.95), Inches(11.2), Inches(2.3),
      [[("* Illustrative estimate, not a guarantee — uses the same benchmarks the app's Metrics tab ships with.", 13.5, NAVY, True, False)],
       [("Basis: a full-system FedRAMP Moderate ATO (~325 controls, ~975 CCIs). Manual A&A benchmark = $700/control "
         "(≈ $233/CCI at ~3 CCIs/control) and 8 hrs/control — the FedRAMP Mod 3PAO range ($150K–$300K, ", 13, INK, False, False),
         ("GAO-24-106591", 13, BLUE, False, False, "https://www.gao.gov/products/GAO-24-106591"),
         (") over ~325 controls, and the industry-standard 8 hr/control per NIST SP 800-53A. Manual A&A "
         "inter-rater agreement runs ~64% (Cohen's κ, ", 13, INK, False, False),
         ("Radziwill & Benton 2017", 13, BLUE, False, False, "https://arxiv.org/abs/1707.02653"),
         ("). Deterministic auto-resolution + AI-assisted assessment remove ~80% of hands-on effort: ≈$180K and "
         "≈2,080 hrs of internal preparation and assessment labor. The platform reduces prep effort, improves "
         "consistency, and lowers assessment costs while supporting assessor activities rather than replacing them.",
         13, INK, False, False)],
       [("Actual savings vary by system size, baseline, evidence quality, and team. Your numbers will differ — these are planning estimates only.",
         13, SLATE, False, True)]], line_spacing=1.12)


# ===========================================================================
# Slide 9 — Proof / close
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_grad_bg(s, NAVY_DEEP, NAVY3, angle=60)
_grad_rect(s, 0, 0, W, Inches(0.14), BLUE_BRIGHT, BLUE, angle=0)
_rect(s, 0, H - Inches(0.1), W, Inches(0.1), BLUE)
_kicker(s, Inches(0.8), Inches(0.6), "The bottom line", SKY)
_text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.9),
      [[("Built for the people who own the ATO", 32, WHITE, True, False)]])

stats = [
    ("7", "frameworks supported"),
    ("27", "evidence file types read"),
    ("5-tier", "evidence-to-control engine"),
    ("U → TS", "every classification enclave"),
]
sx = Inches(0.8)
sw = Inches(2.85)
sgap = Inches(0.13)
for i, (big, small) in enumerate(stats):
    x = sx + (sw + sgap) * i
    _round(s, x, Inches(2.05), sw, Inches(1.5), NAVY2, line=BLUE, line_w=Pt(1), radius=0.07)
    _text(s, x, Inches(2.2), sw, Inches(0.8),
          [[(big, 38, SKY, True, False)]], align=PP_ALIGN.CENTER)
    _text(s, x, Inches(3.0), sw, Inches(0.5),
          [[(small, 12.5, LIGHT, False, False)]], align=PP_ALIGN.CENTER)

_grad_round(s, Inches(0.8), Inches(4.0), Inches(11.73), Inches(1.5), NAVY2, NAVY, angle=60, line=BLUE, line_w=Pt(1.25), radius=0.06)
_grad_rect(s, Inches(0.8), Inches(4.0), Inches(0.16), Inches(1.5), BLUE_BRIGHT, BLUE, angle=90)
_text(s, Inches(1.2), Inches(4.22), Inches(11.0), Inches(1.1),
      [[("Cybersecurity Assessor compresses the slowest, most error-prone phase of authorization into a fast, "
         "repeatable, auditable workflow — your files stay on the workstation, and only the prompt you assess "
         "goes to the AI endpoint you control.",
         15, WHITE, False, False)]], line_spacing=1.18)

_text(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
      [[("Self-contained Windows installer  ·  No separate runtime  ·  Bring your own AI key", 14, SKY, True, False)]])
_text(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.5),
      [[("Cybersecurity Assessor — assess with confidence, defend with evidence.  ·  Nuon", 13, MIST, False, True)]])


# ===========================================================================
# Slide 10 — Competitive differentiation (capability contrast, unnamed)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_grad_bg(s, NAVY_DEEP, NAVY3, angle=60)
_grad_rect(s, 0, 0, W, Inches(0.14), BLUE_BRIGHT, BLUE, angle=0)
_kicker(s, Inches(0.8), Inches(0.45), "How it's different", SKY)
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("Built to assess, not just to track", 32, WHITE, True, False)]])
_text(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.5),
      [[("The work other tools leave to you — finding evidence, mapping it, writing the verdict — this one does for you.",
         15, SKY, False, False)]])

# left card — typical platforms (muted)
lx, ly = Inches(0.8), Inches(2.05)
cardw, cardh = Inches(5.75), Inches(4.55)
_round(s, lx, ly, cardw, cardh, NAVY2, line=SLATE, line_w=Pt(1), radius=0.05)
_grad_rect(s, lx, ly, cardw, Inches(0.09), SLATE, MIST, angle=0)
_text(s, lx + Inches(0.4), ly + Inches(0.28), cardw - Inches(0.8), Inches(0.5),
      [[("Other tools — you do it by hand", 17, MIST, True, False)]])
typical = [
    "Hunt for the right artifact across SharePoint, shares, and inboxes.",
    "Map each artifact to the CCI and control it proves, manually.",
    "Read every scan, STIG, and screenshot and write each verdict yourself.",
    "Hand-type results back into eMASS, then fix the broken formatting.",
    "Build POA&M entries for every gap from scratch.",
    "Re-key citations and hope they survive a 3PAO's scrutiny.",
    "Assemble the Security Assessment Report (SAR) by hand at the end.",
]
_text(s, lx + Inches(0.4), ly + Inches(0.95), cardw - Inches(0.8), Inches(3.3),
      [[("•  " + t, 12, LIGHT, False, False)] for t in typical], line_spacing=1.04, space_after=5)

# right card — Cybersecurity Assessor (highlighted)
rx = lx + cardw + Inches(0.4)
_grad_round(s, rx, ly, cardw, cardh, NAVY, NAVY3, angle=60, line=BLUE, line_w=Pt(1.5), radius=0.05)
_grad_rect(s, rx, ly, cardw, Inches(0.09), BLUE_BRIGHT, BLUE, angle=0)
_text(s, rx + Inches(0.4), ly + Inches(0.28), cardw - Inches(0.8), Inches(0.5),
      [[("Cybersecurity Assessor — done for you", 17, WHITE, True, False)]])
ours = [
    "Sweeps your evidence sources and surfaces the files that matter.",
    "Auto-correlates 27 file types to the exact control they prove.",
    "Proposes a cited verdict for every control — engine-led, AI-assisted.",
    "Writes results back into eMASS in place — formatting preserved.",
    "Generates grounded POA&M entries for every gap automatically.",
    "Cites the exact document, section, and STIG rule — built to defend.",
    "Produces the Security Assessment Report (SAR) when you're done.",
]
_text(s, rx + Inches(0.4), ly + Inches(0.95), cardw - Inches(0.8), Inches(3.4),
      [[("✓  " + t, 12, LIGHT, False, False)] for t in ours], line_spacing=1.04, space_after=5)

# accuracy proof band
_text(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
      [[(">90% verdict accuracy", 14, SKY, True, False),
        ("  ·  validated across 3 real-world system assessments", 14, LIGHT, False, False)]],
      align=PP_ALIGN.CENTER)


# ===========================================================================
# Slide 11 — Why assessors trust it
# ===========================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, WHITE)
_accentbar(s)
_kicker(s, Inches(0.8), Inches(0.45), "Why assessors trust it")
_text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8),
      [[("Defensible by construction, not by claim", 32, NAVY, True, False)]])
_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
      [[("Every design choice answers one question: can the assessor defend this verdict to a reviewer?",
         16, SLATE, False, False)]])

trust = [
    ("Deterministic-first logic",
     "Rules modeled on assessor judgment resolve the clear cases first — repeatable, inspectable, and identical on every run."),
    ("Engine leads, AI assists",
     "Your engine does the assessing and calls the LLM only on genuine judgment calls — then re-validates its answer. Never an unchecked oracle."),
    ("Evidence-backed decisions",
     "Every verdict cites the exact document, section, and STIG rule it rests on — no invented or unsourced citations."),
    ("Full audit trail",
     "Each decision stores the precise evidence snippet behind it, so a reviewer can replay exactly what the assessor saw."),
    ("Human makes the call",
     "Nothing is final until a person accepts it — proposals are editable, and every change is captured for the record."),
    ("Knows its limits — few as they are",
     "Every proposal is confidence-scored; when evidence is thin, it routes the control to a human instead of guessing."),
]
gx, gy = Inches(0.8), Inches(2.2)
gw, gh = Inches(3.84), Inches(1.95)
ggap = Inches(0.1)
for i, (h, b) in enumerate(trust):
    col = i % 3
    rowi = i // 3
    x = gx + (gw + ggap) * col
    y = gy + (gh + Inches(0.18)) * rowi
    _round(s, x, y, gw, gh, SOFT, line=LINE, line_w=Pt(1), radius=0.06)
    _grad_rect(s, x, y, gw, Inches(0.08), BLUE, BLUE_BRIGHT, angle=0)
    _text(s, x + Inches(0.28), y + Inches(0.2), gw - Inches(0.5), Inches(0.7),
          [[(h, 15.5, NAVY, True, False)]], line_spacing=1.0)
    _text(s, x + Inches(0.28), y + Inches(0.92), gw - Inches(0.5), Inches(0.95),
          [[(b, 12, INK, False, False)]], line_spacing=1.1)

# bottom-line banner: the one idea worth keeping from the old "when AI is wrong" slide
_round(s, Inches(0.8), Inches(6.62), Inches(11.73), Inches(0.6), NAVY, line=BLUE, line_w=Pt(1.25), radius=0.18)
_text(s, Inches(0.8), Inches(6.62), Inches(11.73), Inches(0.6),
      [[("Designed so the AI can't be wrong silently: an unsupported proposal degrades to ", 13, LIGHT, False, False),
        ("“needs human review,”", 13, SKY, True, False),
        (" never a wrong verdict written into eMASS.", 13, LIGHT, False, False)]],
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Reorder: the two credibility slides (competitive / trust) are built last
# (indices 9,10) but belong AFTER "Why it's different" (index 6) and BEFORE
# ROI (index 7) + close (index 8). Reorder the sldIdLst so the narrative flows.
# Build order (0-based): 0 Title,1 Problem,2 Enclave,3 How,4 Hood,5 What,
#   6 Why-different,7 ROI,8 Bottom-line,9 Competitive,10 Trust
# Target order:          0,1,2,3,4,5,6, 9,10, 7,8
target = [0, 1, 2, 3, 4, 5, 6, 9, 10, 7, 8]
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for sid in ids:
    sldIdLst.remove(sid)
for i in target:
    sldIdLst.append(ids[i])

out = Path(__file__).resolve().parent / "Cybersecurity_Assessor_Overview.pptx"
prs.save(str(out))
print("WROTE", out, f"({out.stat().st_size // 1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
