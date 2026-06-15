"""PPTX extractor (python-pptx, lazily imported).

Walks every shape on every slide, gathering text from text frames,
tables, and shape titles. Slide notes are appended at the end so they
don't dilute the slide-body context.
"""

from __future__ import annotations

from pathlib import PurePosixPath
from typing import BinaryIO

from ...models import EvidenceKind
from .base import ExtractedDoc, ExtractorError, register, resolve_doc_number


def _shape_text(shape) -> list[str]:
    """Pull text from any shape: text frame, table cell, or grouped shape."""
    out: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            # Use para.text, not a runs-only join. python-pptx exposes text on
            # the paragraph that includes content not carried by explicit runs
            # (e.g. field/auto-text and some templated placeholders); joining
            # runs alone silently drops those lines, gutting the extracted text
            # of evidence decks (matches docx.py, which reads para.text).
            line = para.text.strip()
            if line:
                out.append(line)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                out.append(" | ".join(cells))
    if getattr(shape, "shape_type", None) == 6:  # GROUP
        for sub in shape.shapes:
            out.extend(_shape_text(sub))
    return out


@register(".pptx")
def extract_pptx(stream: BinaryIO, name: str) -> ExtractedDoc:
    """Extract slide + notes text from a PowerPoint deck."""
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        raise ExtractorError(
            "python-pptx is not installed — add it to backend/pyproject.toml "
            "to extract PPTX evidence."
        ) from exc

    stem = PurePosixPath(name).stem

    try:
        prs = Presentation(stream)
    except Exception as exc:  # pragma: no cover
        raise ExtractorError(f"python-pptx failed on {name}: {exc}") from exc

    slide_chunks: list[str] = []
    notes_chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            slide_lines.extend(_shape_text(shape))
        if len(slide_lines) > 1:
            slide_chunks.append("\n".join(slide_lines))
        if slide.has_notes_slide:
            # has_notes_slide can be True while notes_text_frame is None: a
            # notes slide object exists (template/placeholder) but carries no
            # text frame. Guard or python-pptx raises AttributeError on .text
            # and aborts extraction of the whole deck (seen on ACAS scan decks).
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame is not None else ""
            if notes:
                notes_chunks.append(f"## Notes — Slide {i}\n{notes}")

    text = "\n\n".join(slide_chunks + notes_chunks)
    title = (prs.core_properties.title or "").strip() or stem

    return ExtractedDoc(
        text=text,
        title=title,
        doc_number=resolve_doc_number(name, title, text),
        kind=EvidenceKind.PPTX,
        metadata={"slide_count": len(prs.slides)},
    )
