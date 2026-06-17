"""Generate demo evidence artifacts for the Example System Demo system.

Run once to populate demo/policies/, demo/configs/, demo/diagrams/. Idempotent.
"""

from __future__ import annotations

from pathlib import Path

DEMO = Path(__file__).parent
POLICIES = DEMO / "policies"
CONFIGS = DEMO / "configs"
DIAGRAMS = DEMO / "diagrams"


# ---------------------------------------------------------------------------
# DOCX -- Account Management Policy
# ---------------------------------------------------------------------------


def build_account_mgmt_docx() -> Path:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.core_properties.title = "Information System Account Management Policy"
    doc.core_properties.author = "Example System Demo - ISSM"
    doc.core_properties.subject = "USD20240315"

    doc.add_heading("Information System Account Management Policy", level=0)
    p = doc.add_paragraph()
    p.add_run("Document Number: USD20240315").bold = True
    doc.add_paragraph("Version: 2.1")
    doc.add_paragraph("Effective Date: 2024-03-15")
    doc.add_paragraph("Last Reviewed: 2026-04-02")
    doc.add_paragraph("System: Example System Demo (Example System Example System Demo IATT)")

    doc.add_heading("1. Purpose", level=1)
    doc.add_paragraph(
        "This policy establishes the procedures for creating, modifying, "
        "disabling, and removing user accounts on the Example System Demo system in "
        "accordance with NIST SP 800-53 Rev. 5 control AC-2 and supporting "
        "control enhancements AC-2(1) through AC-2(13)."
    )

    doc.add_heading("2. Scope", level=1)
    doc.add_paragraph(
        "All information system accounts on the Example System Demo authorization "
        "boundary, including: individual user accounts, group accounts, "
        "system accounts, application accounts, guest/anonymous accounts, "
        "emergency accounts, temporary accounts, and service accounts."
    )

    doc.add_heading("3. Account Types and Approval Authority", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Account Type"
    hdr[1].text = "Approval Authority"
    hdr[2].text = "Review Interval"
    rows = [
        ("Individual user", "ISSM + supervisor", "Quarterly"),
        ("Privileged (admin)", "ISSM + ISSO", "Monthly"),
        ("Service / application", "ISSM + System Owner", "Quarterly"),
        ("Emergency", "ISSM (after-action)", "Per use"),
        ("Temporary (<= 90 days)", "ISSM + supervisor", "At expiration"),
    ]
    for r in rows:
        row = table.add_row().cells
        for i, v in enumerate(r):
            row[i].text = v

    doc.add_heading("4. Account Lifecycle", level=1)
    for step in [
        "4.1 Request. HR or system owner submits account request ticket.",
        "4.2 Approval. ISSM and supervisor sign-off recorded in ticket.",
        "4.3 Creation. Admin creates account with least-privilege role.",
        "4.4 Notification. User and supervisor notified out-of-band.",
        "4.5 Review. Accounts reviewed per the interval in Section 3.",
        "4.6 Disable. Accounts disabled within 24 hours of departure or "
        "no-longer-needed determination.",
        "4.7 Removal. Disabled accounts retained 90 days then removed.",
    ]:
        doc.add_paragraph(step)

    doc.add_heading("5. Automated Enforcement", level=1)
    doc.add_paragraph(
        "Account lockout, password complexity, password age, and session "
        "termination thresholds are enforced via Active Directory Group "
        "Policy. Current settings are documented in "
        "GPO_Password_Policy_Export.xlsx (USD20240218)."
    )

    doc.add_heading("6. Audit and Logging", level=1)
    doc.add_paragraph(
        "Account management events (Windows Event IDs 4720, 4722, 4724, "
        "4725, 4726, 4738) are forwarded to the Example System SIEM and reviewed "
        "weekly. See AU-6 procedures (USD20240518)."
    )

    doc.add_heading("7. References", level=1)
    for ref in [
        "NIST SP 800-53 Rev. 5: AC-2, AC-2(1), AC-2(3), AC-2(4), AC-2(11), AC-2(12), AC-2(13)",
        "NIST SP 800-53 Rev. 5: IA-5, IA-5(1)",
        "Example System Demo System Security Plan (USD20240101)",
        "Identification and Authentication Procedures (USD20240212)",
    ]:
        doc.add_paragraph(ref, style="List Bullet")

    out = POLICIES / "Information_System_Account_Management_Policy_USD20240315.docx"
    doc.save(str(out))
    return out


# ---------------------------------------------------------------------------
# PDF -- Identification and Authentication Procedures
# ---------------------------------------------------------------------------


def build_ia_procedures_pdf() -> Path:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        PageBreak,
        Table,
        TableStyle,
    )
    from reportlab.lib import colors

    out = POLICIES / "Identification_and_Authentication_Procedures_USD20240212.pdf"
    doc = SimpleDocTemplate(
        str(out),
        pagesize=LETTER,
        leftMargin=0.9 * inch,
        rightMargin=0.9 * inch,
        topMargin=0.9 * inch,
        bottomMargin=0.9 * inch,
        title="Identification and Authentication Procedures",
        author="Example System Demo - ISSM",
        subject="USD20240212",
    )
    styles = getSampleStyleSheet()
    body = styles["BodyText"]
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    title = ParagraphStyle("title", parent=styles["Title"], spaceAfter=12)

    story = []
    story.append(Paragraph("Identification and Authentication Procedures", title))
    story.append(Paragraph("<b>Document Number:</b> USD20240212", body))
    story.append(Paragraph("<b>Version:</b> 1.4", body))
    story.append(Paragraph("<b>Effective Date:</b> 2024-02-12", body))
    story.append(Paragraph("<b>Last Reviewed:</b> 2026-03-19", body))
    story.append(Paragraph("<b>System:</b> Example System Demo (Example System Example System Demo IATT)", body))
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("1. Purpose", h1))
    story.append(
        Paragraph(
            "These procedures implement NIST SP 800-53 Rev. 5 controls IA-2, "
            "IA-3, IA-4, IA-5, IA-6, IA-7, IA-8, and IA-11 on the Example System Demo "
            "system. They define how organizational and non-organizational "
            "users are uniquely identified, authenticated, and their "
            "authenticators managed.",
            body,
        )
    )

    story.append(Paragraph("2. Authentication Methods", h1))
    data = [
        ["User Class", "Primary", "Backup"],
        ["Privileged (admin)", "PIV smart card + PIN", "FIDO2 hardware key"],
        ["Standard user", "Username + password + TOTP", "Recovery code"],
        ["Service account", "Managed Service Account (gMSA)", "n/a"],
        ["Local console", "Username + password", "n/a"],
    ]
    tbl = Table(data, colWidths=[1.6 * inch, 2.5 * inch, 2.1 * inch])
    tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
            ]
        )
    )
    story.append(tbl)
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("3. Password Requirements", h1))
    story.append(
        Paragraph(
            "Password length, complexity, age, history, and lockout settings "
            "are enforced by the Default Domain Policy GPO on the Example System Demo "
            "domain. The current export is recorded in "
            "GPO_Password_Policy_Export.xlsx (USD20240218). Minimum settings:",
            body,
        )
    )
    for item in [
        "Minimum length: 15 characters",
        "Complexity: enabled (3 of 4 character classes)",
        "Maximum age: 60 days",
        "History: 24 passwords remembered",
        "Lockout threshold: 5 invalid attempts in 15 minutes",
        "Lockout duration: 30 minutes; auto-unlock after duration",
    ]:
        story.append(Paragraph("&bull; " + item, body))

    story.append(PageBreak())
    story.append(Paragraph("4. Authenticator Management (IA-5)", h1))
    story.append(
        Paragraph(
            "Initial authenticators are distributed out-of-band via the "
            "issuance station in Room 204. Users must change initial "
            "authenticators on first logon. Lost authenticators are "
            "reported to the help desk; identity is re-verified per "
            "Section 4.2 of USD20240315 before re-issuance.",
            body,
        )
    )

    story.append(Paragraph("4.1 PIV Card Lifecycle", h2))
    story.append(
        Paragraph(
            "PIV cards are issued by the Example System Demo PIV sponsor per FIPS 201-3. "
            "Card termination follows immediately on departure; issuance "
            "and termination events are recorded in the PIV management "
            "system.",
            body,
        )
    )

    story.append(Paragraph("5. Cryptographic Module Authentication (IA-7)", h1))
    story.append(
        Paragraph(
            "All authenticator verification uses FIPS 140-3 validated "
            "cryptographic modules. Microsoft Cryptographic Primitives "
            "Library (Bcryptprimitives.dll) certificate #4544 is in use on "
            "all Windows endpoints; YubiKey 5 FIPS series (certificate "
            "#4569) is used for FIDO2 backup.",
            body,
        )
    )

    story.append(Paragraph("6. References", h1))
    for ref in [
        "NIST SP 800-53 Rev. 5: IA-2, IA-3, IA-4, IA-5, IA-5(1), IA-6, IA-7, IA-8, IA-11",
        "FIPS 201-3 (PIV)",
        "FIPS 140-3 (Cryptographic Module Validation)",
        "NIST SP 800-63B (Digital Identity Guidelines)",
        "Account Management Policy (USD20240315)",
    ]:
        story.append(Paragraph("&bull; " + ref, body))

    doc.build(story)
    return out


# ---------------------------------------------------------------------------
# PPTX -- Security Awareness Training Brief
# ---------------------------------------------------------------------------


def build_training_pptx() -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Example System Demo - Security Awareness Training"
    sub = slide.placeholders[1]
    sub.text = (
        "Q2 2026 Brief | USD20240518 | Annual Refresh\n"
        "AT-2 / AT-3 Implementing Documentation"
    )

    # Helper for bullet slides
    def bullet_slide(title: str, bullets: list[str]) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        body = s.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
            p.level = 0

    bullet_slide(
        "Why You're Here",
        [
            "Annual mandatory training - NIST 800-53 AT-2",
            "Role-based supplement for privileged users - AT-3",
            "Completion tracked in the Example System Demo training roster",
            "Must complete within 30 days of system access grant",
            "Refresher required at least every 12 months",
        ],
    )

    bullet_slide(
        "Insider Threat Awareness (AT-2(2))",
        [
            "Recognize behavioral indicators (unusual hours, data hoarding)",
            "Report concerns via the Example System Demo insider-threat hotline",
            "Never retaliate; reporting is protected",
            "All workstations subject to user-activity monitoring per banner",
        ],
    )

    bullet_slide(
        "Phishing & Social Engineering",
        [
            "Verify sender domain on every external email",
            "Use the Outlook 'Report Phish' button - do not delete",
            "Never enter credentials from a link in email",
            "If you click, report to the help desk within 1 hour",
        ],
    )

    bullet_slide(
        "Handling CUI and Classified Data",
        [
            "Example System Demo is approved up to CUI only",
            "Mark every document with the highest-sensitivity content",
            "Use the approved CUI cover sheet on printed material",
            "Removable media is prohibited without written exception",
        ],
    )

    bullet_slide(
        "Privileged-User Refresher (AT-3)",
        [
            "Use separate accounts for admin work; never browse the web as admin",
            "Log into the privileged access workstation (PAW) for admin tasks",
            "All admin actions are logged - review your own activity weekly",
            "Emergency-access procedure: ISSM approval after the fact within 24 h",
        ],
    )

    bullet_slide(
        "Acknowledgement",
        [
            "Sign the training acknowledgement in the LMS",
            "Roster updates within 24 hours",
            "Questions: ISSM Noah Jaskolski | issm@demo.local",
            "Next refresh due: 12 months from completion",
        ],
    )

    out = POLICIES / "Security_Awareness_Training_Brief_2026Q2_USD20240518.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# XLSX -- GPO Password Policy Export
# ---------------------------------------------------------------------------


def build_gpo_export_xlsx() -> Path:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "Password Policy"

    bold = Font(bold=True)
    header_fill = PatternFill("solid", fgColor="D9E1F2")
    thin = Side(border_style="thin", color="999999")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    wrap = Alignment(wrap_text=True, vertical="top")

    ws["A1"] = "Example System Demo - Default Domain Policy GPO Export"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "Document Number: USD20240218"
    ws["A3"] = "Exported From: example-system-demo-dc01.demo.local"
    ws["A4"] = "Exported By: admin01"
    ws["A5"] = "Export Date: 2026-04-21"
    ws["A6"] = "Policy Path: Computer Configuration > Policies > Windows Settings > Security Settings > Account Policies"
    for r in range(1, 7):
        ws.cell(row=r, column=1).alignment = wrap

    headers = ["Setting", "Value", "NIST 800-53 Control", "Notes"]
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=8, column=c, value=h)
        cell.font = bold
        cell.fill = header_fill
        cell.border = border

    rows = [
        ("Enforce password history", "24 passwords remembered", "IA-5(1)(e)", ""),
        ("Maximum password age", "60 days", "IA-5(1)(d)", ""),
        ("Minimum password age", "1 day", "IA-5(1)", ""),
        ("Minimum password length", "15 characters", "IA-5(1)(a)", "Exceeds DoD 14-char minimum"),
        ("Password must meet complexity requirements", "Enabled", "IA-5(1)(a)", "Requires 3 of 4 char classes"),
        ("Store passwords using reversible encryption", "Disabled", "IA-5(1)(c)", ""),
        ("Account lockout duration", "30 minutes", "AC-7(b)", "Auto-unlock after duration"),
        ("Account lockout threshold", "5 invalid logon attempts", "AC-7(a)", ""),
        ("Reset account lockout counter after", "15 minutes", "AC-7", ""),
        ("Interactive logon: Machine inactivity limit", "900 seconds", "AC-11(a)", "15-minute screen lock"),
        ("Microsoft network server: Idle session", "15 minutes", "AC-12", ""),
        ("Network security: Force logoff when logon hours expire", "Enabled", "AC-2(11)", ""),
    ]
    for ri, (setting, value, ctl, notes) in enumerate(rows, start=9):
        ws.cell(row=ri, column=1, value=setting).border = border
        ws.cell(row=ri, column=2, value=value).border = border
        ws.cell(row=ri, column=3, value=ctl).border = border
        ws.cell(row=ri, column=4, value=notes).border = border
        for c in range(1, 5):
            ws.cell(row=ri, column=c).alignment = wrap

    ws.column_dimensions["A"].width = 52
    ws.column_dimensions["B"].width = 32
    ws.column_dimensions["C"].width = 22
    ws.column_dimensions["D"].width = 40

    # Second sheet: account lockout & session settings cross-reference
    ws2 = wb.create_sheet("Control Cross-Reference")
    ws2["A1"] = "Settings -> NIST 800-53 Control Mapping"
    ws2["A1"].font = Font(bold=True, size=13)
    mapping_hdr = ["Control", "Settings That Implement It"]
    for c, h in enumerate(mapping_hdr, start=1):
        cell = ws2.cell(row=3, column=c, value=h)
        cell.font = bold
        cell.fill = header_fill
        cell.border = border
    mapping = [
        ("AC-2(11)", "Force logoff when logon hours expire"),
        ("AC-7", "Account lockout duration, threshold, reset"),
        ("AC-11", "Machine inactivity limit (screen lock)"),
        ("AC-12", "Idle session disconnect"),
        ("IA-5(1)", "Password history/age/length/complexity/encryption"),
    ]
    for ri, (ctl, setting) in enumerate(mapping, start=4):
        ws2.cell(row=ri, column=1, value=ctl).border = border
        ws2.cell(row=ri, column=2, value=setting).border = border
        for c in range(1, 3):
            ws2.cell(row=ri, column=c).alignment = wrap
    ws2.column_dimensions["A"].width = 16
    ws2.column_dimensions["B"].width = 60

    out = CONFIGS / "GPO_Password_Policy_Export_USD20240218.xlsx"
    wb.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Boundary diagrams — TWO tenant-distinct authorization boundaries
# ---------------------------------------------------------------------------
#
# Example System Demo is a multi-boundary program: customer workloads run in
# BOTH AWS GovCloud and Azure Government (mirroring the two CRMs under crm/).
# Each cloud has its OWN authorization-boundary diagram with cloud-native
# terminology and DISTINCT topology — exactly how a real SSP package ships one
# boundary diagram per enclave. The two diagrams are deliberately NOT
# interchangeable: the AWS one speaks VPC / Security Group / GovCloud, the Azure
# one speaks VNet / NSG / Azure Government. That difference is what makes the
# multi-tenant boundary-attribution story real — each artifact belongs to one
# tenant, and a per-scope (narratives_by_scope) verdict must cite the right one.
#
# The diagram extractor (evidence/extractors/diagram.py) reads the shape/label
# text verbatim, so the labels below are exactly what reaches the LLM. The
# filename + boundary keywords ("boundary", "firewall"/"NSG", subnets) drive the
# tagger's diagram→boundary-control rule (SC-7 / CA-3 / AC-4 / PL-8).
#
# Provided in TWO on-disk formats so both extractor paths are exercised end to
# end: the AWS boundary ships as Visio (.vsdx), the Azure boundary as SVG.

# AWS GovCloud authorization boundary — Visio (.vsdx).
_AWS_BOUNDARY_SHAPES = [
    "Internet / External Users",
    "AWS GovCloud (US-East) Region",
    "Internet Gateway + AWS WAF",
    "Application Load Balancer (public subnet)",
    "Perimeter Security Group - default deny inbound",
    "Production VPC 10.20.0.0/16",
    "Public Subnet (DMZ) 10.20.1.0/24",
    "App Tier - EC2 Auto Scaling Group 10.20.5.0/24",
    "Data Tier - RDS PostgreSQL (Multi-AZ) 10.20.9.0/24",
    "Management Subnet - bastion / SSM 10.20.99.0/24",
    "VPC Flow Logs -> CloudWatch -> Splunk forwarder",
    "Amazon GuardDuty + AWS Config",
    "Authorization Boundary - Example System Demo (AWS GovCloud, IATT)",
]

# Azure Government authorization boundary — SVG. DISTINCT topology + addressing.
_AZURE_BOUNDARY_LABELS = [
    "Internet / External Users",
    "Azure Government (USGov Virginia)",
    "Azure Front Door + Web Application Firewall",
    "Hub VNet 172.16.0.0/16",
    "Perimeter NSG - default deny inbound",
    "DMZ Subnet 172.16.1.0/24",
    "Spoke VNet - App Subnet 172.16.5.0/24",
    "Spoke VNet - Data Subnet (Azure SQL MI) 172.16.9.0/24",
    "Azure Bastion - Management Subnet 172.16.99.0/24",
    "NSG Flow Logs -> Log Analytics -> Microsoft Sentinel",
    "Microsoft Defender for Cloud",
    "Authorization Boundary - Example System Demo (Azure Government, IATT)",
]


def _write_vsdx(out: Path, page_name: str, shapes: list[str]) -> Path:
    """Write a minimal but valid OOXML/vsdx zip whose page carries shape text."""
    import zipfile

    shape_xml = "\n".join(
        f'  <Shape ID="{i}" Type="Shape"><Text>{label}</Text></Shape>'
        for i, label in enumerate(shapes, start=1)
    )
    page_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<PageContents xmlns="http://schemas.microsoft.com/office/visio/2012/main"'
        ' xml:space="preserve">\n'
        " <Shapes>\n"
        f"{shape_xml}\n"
        " </Shapes>\n"
        "</PageContents>"
    )
    pages_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Pages xmlns="http://schemas.microsoft.com/office/visio/2012/main">'
        f'<Page ID="0" Name="{page_name}"/></Pages>'
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.microsoft.com/visio/2010/relationships/document" '
        'Target="visio/document.xml"/></Relationships>'
    )
    DIAGRAMS.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("visio/pages/pages.xml", pages_xml)
        z.writestr("visio/pages/page1.xml", page_xml)
    return out


def build_aws_boundary_diagram_vsdx() -> Path:
    """AWS GovCloud authorization boundary — Visio (.vsdx)."""
    out = DIAGRAMS / "Example_System_AWS_GovCloud_Boundary_Diagram_USD20240620.vsdx"
    return _write_vsdx(out, "AWS GovCloud Authorization Boundary", _AWS_BOUNDARY_SHAPES)


def build_azure_boundary_diagram_svg() -> Path:
    """Azure Government authorization boundary — SVG."""
    text_nodes = "\n".join(
        f'  <text x="40" y="{50 + i * 42}">{label}</text>'
        for i, label in enumerate(_AZURE_BOUNDARY_LABELS)
    )
    svg = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<svg xmlns="http://www.w3.org/2000/svg" width="680" height="600">\n'
        "  <title>Example System Demo - Azure Government Authorization Boundary</title>\n"
        "  <desc>Azure Government (USGov Virginia) network boundary and "
        "segmentation diagram</desc>\n"
        f"{text_nodes}\n"
        "</svg>"
    )
    DIAGRAMS.mkdir(parents=True, exist_ok=True)
    out = DIAGRAMS / "Example_System_Azure_Government_Boundary_Diagram_USD20240621.svg"
    out.write_text(svg, encoding="utf-8")
    return out


# ---------------------------------------------------------------------------
# DOCX -- AWS GovCloud Remote Access configuration (AC-17 customer-side)
# ---------------------------------------------------------------------------
#
# Evidence for the AWS GovCloud CUSTOMER scope of AC-17. The Azure scope is
# inherited (managed Azure Bastion) and needs no customer artifact, but the AWS
# scope is customer-owned, so the assessor needs a real artifact to assess that
# half — without it the control correctly abstains. This doc cites "AC-17" and
# "CCI-000063" in body text so the tagger's Tier-3 (control-ID) and CCI passes
# map it to the right objective.


def build_remote_access_config_docx() -> Path:
    from docx import Document

    doc = Document()
    doc.core_properties.title = "AWS GovCloud Remote Access Configuration"
    doc.core_properties.author = "Example System Demo - Cloud Engineering"
    doc.core_properties.subject = "USD20240622"

    doc.add_heading("AWS GovCloud Remote Access Configuration", level=0)
    p = doc.add_paragraph()
    p.add_run("Document Number: USD20240622").bold = True
    doc.add_paragraph("Version: 1.0")
    doc.add_paragraph("Effective Date: 2026-05-01")
    doc.add_paragraph("Scope: AWS GovCloud (US) enclave of the Example System Demo")

    doc.add_heading("1. Purpose", level=1)
    doc.add_paragraph(
        "This document records the customer-configured remote-access controls "
        "for the AWS GovCloud enclave, satisfying the customer responsibility "
        "for NIST SP 800-53 Rev. 5 control AC-17 (Remote Access), CCI-000063, "
        "as assigned to the customer in the AWS GovCloud Customer Responsibility "
        "Matrix. (On the Azure Government enclave AC-17 is inherited via managed "
        "Azure Bastion and requires no customer configuration.)"
    )

    doc.add_heading("2. Remote Access Method", level=1)
    doc.add_paragraph(
        "All administrative remote access to AWS GovCloud workloads is brokered "
        "through AWS Client VPN with mutual-certificate authentication federated "
        "to the enterprise IdP. No customer EC2 instance exposes a public RDP or "
        "SSH port; security groups deny inbound administrative traffic from the "
        "internet."
    )

    doc.add_heading("3. Authorization and Enforcement", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Setting"
    hdr[1].text = "Configured Value"
    for setting, value in [
        ("Remote-access method", "AWS Client VPN (mutual TLS + IdP)"),
        ("Authorization prior to connection", "Required — conditional access policy"),
        ("Multi-factor authentication", "Enforced for all VPN sessions"),
        ("Split tunneling", "Disabled"),
        ("Session logging", "VPC Flow Logs + CloudWatch, 1-year retention"),
        ("Entitlement review", "Quarterly"),
    ]:
        row = table.add_row().cells
        row[0].text = setting
        row[1].text = value

    doc.add_heading("4. Assessment Notes", level=1)
    doc.add_paragraph(
        "Remote access is authorized prior to connection and enforces approved "
        "methods per AC-17. Evidence: AWS Client VPN endpoint configuration "
        "export and conditional-access policy, retained in the Example System "
        "Demo cloud-engineering record set under document USD20240622."
    )

    out = POLICIES / "AWS_GovCloud_Remote_Access_Configuration_USD20240622.docx"
    doc.save(str(out))
    return out


# ---------------------------------------------------------------------------
# DOCX -- Audit Log Storage Capacity Memo (AU-4) — DELIBERATELY CONTRADICTORY
# ---------------------------------------------------------------------------
#
# Evidence for the AU-4 abstain showcase. This artifact INTERNALLY CONTRADICTS
# itself: Section 2 asserts the allocated capacity meets the 1-year retention
# requirement, while Section 4 records a utilization analysis showing the same
# partition fills in ~40 days and the expansion is only PLANNED. A control whose
# only evidence both affirms and refutes the requirement is the textbook case
# the assessor must NOT auto-rule: the LLM legitimately cannot reach a confident
# verdict, so it abstains (needs_review) and a human adjudicates. The contradiction
# is what drives low-confidence / dual-pass disagreement → abstain, rather than a
# fabricated Compliant or Non-Compliant. Cites "AU-4" + "CCI-000137" in body text
# so the tagger maps it to the right objective.


def build_audit_storage_memo_docx() -> Path:
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Audit Log Storage Capacity Memo"
    doc.core_properties.author = "Example System Demo - System Engineering"
    doc.core_properties.subject = "USD20240623"

    doc.add_heading("Audit Log Storage Capacity Memo", level=0)
    p = doc.add_paragraph()
    p.add_run("Document Number: USD20240623").bold = True
    doc.add_paragraph("Version: 0.9 (DRAFT — pending engineering sign-off)")
    doc.add_paragraph("Effective Date: 2026-05-10")
    doc.add_paragraph("System: Example System Demo (Example System Example System Demo IATT)")

    doc.add_heading("1. Purpose", level=1)
    doc.add_paragraph(
        "This memo records the audit log storage capacity allocation for the "
        "Example System Demo SIEM tier, addressing NIST SP 800-53 Rev. 5 control "
        "AU-4 (Audit Log Storage Capacity), CCI-000137. The organization-defined "
        "requirement is to retain audit records online for one (1) year without "
        "loss of records."
    )

    doc.add_heading("2. Allocated Capacity", level=1)
    doc.add_paragraph(
        "A dedicated 500 GB partition is allocated to the SIEM indexer for audit "
        "record storage. Based on the original sizing estimate, this allocation "
        "is sufficient to meet the 1-year online retention requirement, and AU-4 "
        "is therefore satisfied."
    )

    doc.add_heading("3. Monitoring", level=1)
    doc.add_paragraph(
        "Storage utilization is monitored via a CloudWatch alarm that pages the "
        "on-call engineer at 80% capacity."
    )

    doc.add_heading("4. Utilization Analysis (2026-05)", level=1)
    doc.add_paragraph(
        "NOTE: A utilization analysis completed 2026-05-08 found that, at the "
        "current measured ingest rate, the 500 GB partition reaches capacity in "
        "approximately 40 days — far short of the 1-year requirement. Expansion "
        "to a 4 TB partition (sized for 1-year retention) is PLANNED for the next "
        "maintenance window but is NOT yet implemented. Until expansion, the "
        "system relies on oldest-record rollover, which would lose audit records "
        "before the 1-year retention period elapses."
    )

    doc.add_heading("5. Status", level=1)
    doc.add_paragraph(
        "This memo is a DRAFT pending engineering sign-off; Sections 2 and 4 have "
        "not yet been reconciled. The capacity posture for AU-4 is unresolved as "
        "of this writing."
    )

    out = CONFIGS / "Audit_Log_Storage_Capacity_Memo_USD20240623.docx"
    doc.save(str(out))
    return out


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------


if __name__ == "__main__":
    for fn in (
        build_account_mgmt_docx,
        build_ia_procedures_pdf,
        build_training_pptx,
        build_gpo_export_xlsx,
        build_remote_access_config_docx,
        build_audit_storage_memo_docx,
        build_aws_boundary_diagram_vsdx,
        build_azure_boundary_diagram_svg,
    ):
        path = fn()
        print(f"WROTE  {path.relative_to(DEMO)}")
